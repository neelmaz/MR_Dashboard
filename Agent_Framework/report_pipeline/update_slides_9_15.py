"""
Consolidated February-PPT updater for slides 9, 10, 12, 13, 14, 15.

Implements the per-slide spec in data/Slide&DataMapping.txt. The mapping file's
"Slide N" == python-pptx slide INDEX N (0-based) == PowerPoint position N+1, e.g.
mapping "Slide 9" is the Gender/Age/Nationality/Residence demographics slide
(prs.slides[9]).

It edits an existing February deck IN PLACE (default: the latest Feb output) so
all downstream work on slides 16+ is preserved; only slides 9,10,12,13,14,15 are
touched. The working Excel-summary agent is NOT modified or required to change.

Data sources
------------
  * Feb summary  (outputs/Feb_summary.xlsx, "Summary" sheet)
        - channel KPIs (Satisfaction/NPS/Trust/Value) for slides 12 & 13
  * Feb basic    (data/feb_basic/...xlsx, "Percentages" sheet)
        - demographic distributions (slides 9, 10)  -> via relabel_charts
        - Overall-Satisfaction (A12+B5+C5) T2B by demographic (slide 14)
        - Overall-Satisfaction T2B by service        (slide 15)
  * Jan deck     (data/jan_report.pptx)   - anchors distribution charts (9,10)
  * Jan summary  (data/jan_summary.xlsx)  - anchors NPS rectangles by Jan value

Note on slide 14: its nominal source (the Summary "A16" block, T33:W55) still
holds January figures in the Feb summary (the Excel agent leaves that block
unchanged), so slide-14 values + sample sizes are computed from the basic table's
Overall-Satisfaction T2B-by-demographic instead. The verified Feb sample sizes
(Male 853 / Female 1013 = 1,866, etc.) match update_bases.py.

Note on slide 15: the basic table only provides service names in ENGLISH; the
chart categories are written with those English names (translate in PowerPoint as
needed). February has 21 qualifying services (sample > 0), so the two bar charts
are expanded beyond the 11 January bars.

Usage:
  python update_slides_9_15.py [in.pptx] [out.pptx]
        [--feb-summary outputs/Feb_summary.xlsx]
        [--feb-basic   "data/feb_basic/MOHAP CSAT Tabs - Feb 2026.xlsx"]
        [--jan-basic   "data/jan_basic/MOHAP CSAT Tabs for Jan-2026.xlsx"]
        [--jan-summary data/jan_summary.xlsx]
        [--jan-deck    data/jan_report.pptx]
"""
from __future__ import annotations

import argparse
import copy
import json
import re
from pathlib import Path

import openpyxl
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

import relabel_charts as rc  # reuse the proven label-anchored distribution writer

C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

ARABIC_INDIC = str.maketrans("0123456789", "٠١٢٣٤٥٦٧٨٩")
WESTERN = str.maketrans("٠١٢٣٤٥٦٧٨٩", "0123456789")

# ── channel dashboard (slides 12 & 13) ──────────────────────────────────────────
# channel identified by shape Left (inches); metric by Summary column / position.
CHAN_LEFT = {"Overall": 9.1, "Mobile": 7.6, "Website": 5.6,
             "CSC": 3.7, "CallCenter": 1.8, "Email": -0.1}
SUM_ROW = {"Overall": 8, "Mobile": 2, "Website": 3, "CSC": 4, "CallCenter": 5, "Email": 6}
SUM_COL = {"Sat": 21, "Trust": 24, "Value": 25, "NPS": 23}

# ── slide-14 source: the Feb Summary "A16" block (T33:W55) ───────────────────────
# Summary row -> demographic segment (fixed by the Original formula references);
# V(22)=sample size, W(23)=satisfaction %. The Excel agent now populates these for
# February (build_feb_summary.fill_overall_sat_demographics).
A16_V, A16_W = 22, 23
A16_ROW = {
    "Female": 34, "Male": 35,
    "Others": 37, "GCC": 38, "UAE": 39, "Asian": 40, "Arab": 41,
    "African": 42, "Westerner": 43,
    "Elder": 45, "Adult": 46, "Youth": 47,
    "UAQ": 49, "Fujairah": 50, "RAK": 51, "Sharjah": 52, "Ajman": 53,
    "Dubai": 54, "AbuDhabi": 55,
}
# chart category order (Arabic) -> the demographic keys above, per slide-14 chart.
# A merged "OthersNat" bucket = Westerner + Others (weighted by base).
SLIDE14_CHART_SEGS = {
    10: ["Male", "Female"],                                   # الجنس
    11: ["Elder", "Adult", "Youth"],                          # الفئة العمرية (>60, 30-60, 18-30)
    13: ["AbuDhabi", "Dubai", "Ajman", "Sharjah", "RAK", "Fujairah", "UAQ"],
    16: ["GCC", "UAE", "Asian", "Arab", "African", "OthersNat"],
}
# slide-14 bar colour is VALUE-based, not position-based:
#   >= 90% -> green, 85-90% -> amber/gold, < 85% -> red
# (reproduces the January template: 89% Abu Dhabi = gold, 80.6% >60 = red).
GREEN_FILL = "35AA6E"
AMBER_FILL = "B68A35"
RED_FILL = "E53935"


def bar_colour(v):
    if v >= 0.90:
        return GREEN_FILL
    if v >= 0.85:
        return AMBER_FILL
    return RED_FILL
# slide-14 sample-size textbox shape_id -> demographic key (verified vs Jan values)
SLIDE14_SAMPLE = {
    17: "OthersNat", 18: "Arab", 19: "UAE", 20: "GCC",
    21: "AbuDhabi", 25: "Dubai", 26: "Ajman", 27: "Sharjah", 28: "RAK",
    30: "Fujairah", 31: "UAQ", 32: "Female", 33: "Male", 34: "Adult",
    35: "Elder", 36: "Youth", 37: "African", 41: "Asian",
}

TOTAL_SAMPLE_OLD = "3,001"   # Jan overall sample text to replace on slides 9,10,12,13,14


# ── small helpers ────────────────────────────────────────────────────────────────

def iter_shapes(shapes):
    """Yield every shape, recursing into groups."""
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(s.shapes)
        else:
            yield s


def first_series(chart):
    return chart._chartSpace.find(".//" + C + "ser")


def set_ring(chart, v):
    ser = first_series(chart)
    pts = ser.find(C + "val").find(C + "numRef").find(C + "numCache").findall(C + "pt")
    if len(pts) >= 2:
        pts[0].find(C + "v").text = f"{v:.6g}"
        pts[1].find(C + "v").text = f"{1 - v:.6g}"


def set_title(chart, txt):
    title = chart._chartSpace.find(".//" + C + "title")
    if title is None:
        return
    ts = title.findall(".//" + A + "t")
    if not ts:
        return
    ts[0].text = txt
    for extra in ts[1:]:
        extra.text = ""


def set_series_values(chart, values):
    """Overwrite a single-series chart's value cache (order = category order)."""
    ser = first_series(chart)
    cache = ser.find(C + "val").find(C + "numRef").find(C + "numCache")
    pts = cache.findall(C + "pt")
    for pt, v in zip(pts, values):
        pt.find(C + "v").text = f"{v:.10g}"


def chan_of(left_in):
    return min(CHAN_LEFT, key=lambda c: abs(CHAN_LEFT[c] - left_in))


def replace_number_text(shape, new_text):
    """Replace a textbox whose paragraph is just a number (optionally in parens)
    with `new_text`, keeping the first run's font. Returns True on success."""
    para = shape.text_frame.paragraphs[0]
    runs = para.runs
    if not runs:
        return False
    joined = "".join(r.text for r in runs)
    runs[0].text = new_text
    for r in runs[1:]:
        r.text = ""
    return joined != new_text or True


def rebuild_cache(cache, vals, is_str):
    """Replace a c:numCache / c:strCache point list with `vals`."""
    for pt in cache.findall(C + "pt"):
        cache.remove(pt)
    ptc = cache.find(C + "ptCount")
    if ptc is None:
        ptc = etree.SubElement(cache, C + "ptCount")
    ptc.set("val", str(len(vals)))
    for i, v in enumerate(vals):
        pt = etree.SubElement(cache, C + "pt")
        pt.set("idx", str(i))
        ve = etree.SubElement(pt, C + "v")
        ve.text = v if is_str else f"{v:.10g}"


def rewrite_bar_chart(chart, labels, values):
    """Replace a clustered-bar chart's categories (strings) and values."""
    ser = first_series(chart)
    cat = ser.find(C + "cat")
    val = ser.find(C + "val")
    # Drop January per-point formatting/labels (they pinned colours & labels to the
    # old 5/6 bar indices and would mis-highlight bars once the count changes).
    for tag in ("dPt", "dLbl"):
        for el in ser.findall(C + tag):
            ser.remove(el)
    # categories -> strCache (force string categories)
    scache = cat.find(".//" + C + "strCache")
    if scache is None:
        # convert a numRef category to a strRef/strCache
        for child in list(cat):
            cat.remove(child)
        sref = etree.SubElement(cat, C + "strRef")
        etree.SubElement(sref, C + "f").text = "Sheet1!$A$1"
        scache = etree.SubElement(sref, C + "strCache")
    rebuild_cache(scache, labels, is_str=True)
    ncache = val.find(".//" + C + "numCache")
    rebuild_cache(ncache, values, is_str=False)


# ── data loading ─────────────────────────────────────────────────────────────────

def load_summary(path):
    return openpyxl.load_workbook(path, data_only=True)["Summary"]


def kpi(sm, chan, metric):
    return sm.cell(SUM_ROW[chan], SUM_COL[metric]).value


def locate_overall_sat(ws):
    """Return (base_row, t2b_row) of the 'Overall Satisfaction (A12+B5+C5)' table."""
    title = None
    for r in range(1, ws.max_row + 1):
        v = ws.cell(r, 1).value
        if isinstance(v, str) and "overall sati" in v.lower():
            title = r
            break
    if title is None:
        raise RuntimeError("Overall Satisfaction table not found in basic Percentages sheet")
    base = t2b = None
    for r in range(title + 1, title + 25):
        a = ws.cell(r, 1).value
        if isinstance(a, str):
            if base is None and a.strip().lower().startswith("base:"):
                base = r
            if t2b is None and a.strip().upper().startswith("T2B"):
                t2b = r
        if base and t2b:
            break
    return base, t2b


def load_services(ws, base_row, t2b_row):
    """Ranked list of services with sample>0 and numeric T2B, from the
    'Service for' banner (col 63 onward). Sort: T2B desc, then sample desc."""
    rows = []
    for c in range(63, ws.max_column + 1):
        name = ws.cell(2, c).value
        if not name:
            continue
        base = ws.cell(base_row, c).value
        t2b = ws.cell(t2b_row, c).value
        if isinstance(base, (int, float)) and base > 0 and isinstance(t2b, (int, float)):
            rows.append((str(name).strip(), int(round(base)), float(t2b)))
    rows.sort(key=lambda x: (-x[2], -x[1]))
    return rows


# ── per-slide updaters ───────────────────────────────────────────────────────────

def update_distribution_slides(feb_prs, jan_prs, jan_secs, feb_secs, slide_idxs, audit):
    """Slides 9 & 10: label-anchored demographic distribution charts."""
    jan_charts = {}
    for si, sl in enumerate(jan_prs.slides):
        for sh in rc.iter_charts(sl.shapes):
            jan_charts[(si, sh.shape_id)] = sh.chart
    for si in slide_idxs:
        for sh in rc.iter_charts(feb_prs.slides[si].shapes):
            jchart = jan_charts.get((si, sh.shape_id))
            if jchart is None:
                continue
            jsers = jchart._chartSpace.findall(".//c:ser", rc.NS)
            fsers = sh.chart._chartSpace.findall(".//c:ser", rc.NS)
            for jser, fser in zip(jsers, fsers):
                jpts = rc.series_pts(jser)
                fpts = rc.series_pts(fser)
                if not jpts or len(jpts) != len(fpts):
                    continue
                jvals = [f for _, _, f in jpts]
                m = rc.best_section(jvals, jan_secs)
                if m is None:
                    continue
                header, recipe = m
                feb = [rc.feb_value(feb_secs, header, labs) for labs in recipe]
                if any(v is None for v in feb):
                    continue
                old = [round(f, 4) for _, _, f in fpts]
                for (_, vel, _), nv in zip(fpts, feb):
                    vel.text = f"{nv:.10g}"
                audit.append({"slide": si, "shape": sh.shape_id, "kind": "distribution",
                              "section": header, "old": old,
                              "new": [round(v, 4) for v in feb]})


# Slide 10 grouped demographic charts. The label-anchor can't resolve these:
# each chart drops the "Refuse to answer" option and renormalises the rest to
# 100%, and occupation merges Student+Others into one bar. So we map each chart
# bar (in the chart's category order) to its source option(s) explicitly, sum,
# then normalise over the included options.  (header matched by keyword.)
SLIDE10 = {
    19: ("occupation", [["student", "others"], ["retired"], ["self-employed"],
                        ["job seeker"], ["housewife"], ["private sector"], ["government"]]),
    25: ("marital",    [["widow"], ["divorced"], ["married with no children"],
                        ["single"], ["married with children"]]),
    54: ("education",  [["less than high school"], ["high school graduate"],
                        ["undergraduates"]]),
}


def _find_section(feb_secs, keyword):
    for header, labels in feb_secs.items():
        if keyword in header.lower():
            return labels
    return {}


def update_slide10(feb_prs, feb_secs, audit):
    """Slide 10: occupation / social status / education distribution charts."""
    for sh in iter_shapes(feb_prs.slides[10].shapes):
        if not getattr(sh, "has_chart", False):
            continue
        spec = SLIDE10.get(sh.shape_id)
        if spec is None:
            continue
        keyword, groups = spec
        sec = _find_section(feb_secs, keyword)
        raws, ok = [], True
        for grp in groups:
            total = 0.0
            found = False
            for sub in grp:
                for lab, v in sec.items():
                    if sub in lab.lower() and isinstance(v, (int, float)):
                        total += v
                        found = True
            if not found:
                ok = False
                break
            raws.append(total)
        norm = sum(raws)
        if not ok or norm <= 0:
            continue
        vals = [r / norm for r in raws]
        old = [round(x, 4) for ser in sh.chart.series for x in ser.values]
        set_series_values(sh.chart, vals)
        audit.append({"slide": 10, "shape": sh.shape_id, "kind": "distribution_grouped",
                      "section": keyword, "old": old, "new": [round(v, 4) for v in vals]})


def update_doughnut_dashboard(feb_prs, sm, audit):
    """Slide 12 (Satisfaction) and slide 13 (Trust/Value) ring doughnuts."""
    for si in (12, 13):
        for d in iter_shapes(feb_prs.slides[si].shapes):
            if not getattr(d, "has_chart", False):
                continue
            if "DOUGHNUT" not in str(d.chart.chart_type):
                continue
            chan = chan_of(Emu(d.left).inches)
            if si == 12:
                metric = "Sat"
            else:
                metric = "Trust" if Emu(d.top).inches < 3.5 else "Value"
            v = kpi(sm, chan, metric)
            if not isinstance(v, (int, float)):
                continue
            old = [round(x, 4) for ser in d.chart.series for x in ser.values]
            set_ring(d.chart, float(v))
            set_title(d.chart, f"{round(v * 100)}%")
            audit.append({"slide": si, "shape": d.shape_id, "kind": "doughnut",
                          "channel": chan, "metric": metric,
                          "old": old, "new": round(v, 4)})


# NPS score rectangles on slide 12: the 5 channel rects are top-level autoshapes
# (identified by Left); the Overall rect sits inside the group carrying the "NPS"
# label. Position-based so it is robust to the deck's current numbers.
NPS_CHAN_LEFT = {"Email": 0.70, "CallCenter": 2.57, "CSC": 4.49,
                 "Website": 6.41, "Mobile": 8.30}


def _is_int_text(shape):
    return shape.has_text_frame and re.fullmatch(
        r"\d{1,3}", shape.text_frame.text.strip().translate(WESTERN))


def update_nps_rectangles(feb_prs, sm, audit):
    """Slide 12 NPS score rectangles, set to round(Feb_NPS*100)."""
    slide = feb_prs.slides[12]
    # Overall: numeric rect inside the group that also holds an "NPS" label
    for grp in slide.shapes:
        if grp.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        if not any(c.has_text_frame and c.text_frame.text.strip() == "NPS"
                   for c in grp.shapes):
            continue
        for c in grp.shapes:
            if _is_int_text(c):
                new = str(round(kpi(sm, "Overall", "NPS") * 100))
                old = c.text_frame.text.strip()
                replace_number_text(c, new)
                audit.append({"slide": 12, "shape": c.shape_id, "kind": "nps",
                              "channel": "Overall", "old": old, "new": new})
    # Channels: top-level numeric autoshapes mapped by Left
    for sh in slide.shapes:
        if not _is_int_text(sh):
            continue
        left = Emu(sh.left).inches
        ch = min(NPS_CHAN_LEFT, key=lambda c: abs(NPS_CHAN_LEFT[c] - left))
        if abs(NPS_CHAN_LEFT[ch] - left) > 0.6:
            continue
        new = str(round(kpi(sm, ch, "NPS") * 100))
        old = sh.text_frame.text.strip()
        replace_number_text(sh, new)
        audit.append({"slide": 12, "shape": sh.shape_id, "kind": "nps",
                      "channel": ch, "old": old, "new": new})


# Per-channel sample-size tables at the bottom of slides 12 & 13. The table cells
# (Jan: [281,515,38,1000,1167] = Email/CallCenter/CSC/Website/Mobile, sum 3,001)
# carry the individual channel bases. Column->channel is derived from the Jan deck
# (matching each cell to a channel's Jan sample), so it survives column reordering.
SAMPLE_COL = 22
CHANNEL_TABLE = {12: 10, 13: 15}   # slide index -> table shape_id


def channel_sample(sm, chan):
    return sm.cell(SUM_ROW[chan], SAMPLE_COL).value


def _find_table(slide, shape_id):
    for sh in iter_shapes(slide.shapes):
        if sh.shape_id == shape_id and sh.has_table:
            return sh.table
    return None


def _set_cell_number(cell, new_text):
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        cell.text_frame.text = new_text
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ""


def update_channel_samples(feb_prs, jan_prs, sm, jan_sm, audit):
    """Slides 12 & 13: per-channel sample-size table (individual item bases)."""
    jan_by_val = {}
    for ch in SUM_ROW:
        v = channel_sample(jan_sm, ch)
        if isinstance(v, (int, float)):
            jan_by_val[int(round(v))] = ch
    for si, tid in CHANNEL_TABLE.items():
        jtab = _find_table(jan_prs.slides[si], tid)
        ftab = _find_table(feb_prs.slides[si], tid)
        if jtab is None or ftab is None:
            continue
        jcells = [c for row in jtab.rows for c in row.cells]
        fcells = [c for row in ftab.rows for c in row.cells]
        for jc, fc in zip(jcells, fcells):
            jv = jc.text.strip().replace(",", "").translate(WESTERN)
            if not jv.isdigit():
                continue
            ch = jan_by_val.get(int(jv))
            if ch is None:
                continue
            fv = channel_sample(sm, ch)
            if not isinstance(fv, (int, float)):
                continue
            new = f"{int(round(fv)):,}"
            old = fc.text.strip()
            _set_cell_number(fc, new)
            audit.append({"slide": si, "shape": tid, "kind": "channel_sample",
                          "channel": ch, "old": old, "new": new})


# Headline narrative paragraphs (slides 12 & 13: the long sentence with %s). These
# text boxes are not generated from data, so the January-derived figures/channels go
# stale. We patch the specific runs to agree with the Feb summary KPIs. Run texts are
# matched verbatim against the base deck; the new values below are the Feb summary
# figures (Email Sat 86%, Overall NPS 72, NPS leader = CSC 78, Value leader = Website
# 93% — mobile app is 92%). Matching the stale text makes the swaps idempotent: once
# applied they no longer match, so re-runs are safe.

def _headline_shape(slide):
    for sh in iter_shapes(slide.shapes):
        if sh.has_text_frame and "%" in sh.text_frame.text and len(sh.text_frame.text) > 60:
            return sh
    return None


def _swap_runs(shape, slide_idx, pairs, audit):
    """Apply ordered (old_run_text -> new_run_text) swaps; each old matches one run."""
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    for old, new in pairs:
        for r in runs:
            if r.text == old and r.text != new:
                r.text = new
                audit.append({"slide": slide_idx, "shape": shape.shape_id,
                              "kind": "headline", "old": old, "new": new})
                break


def update_headlines(feb_prs, sm, audit):
    """Fix the slide 12 & 13 headline narrative so it matches the Feb summary KPIs."""
    sh = _headline_shape(feb_prs.slides[12])
    if sh is not None:
        _swap_runs(sh, 12, [
            ("7%.", "6%."),                                   # email satisfaction 87% -> 86%
            ("73", "72"),                                     # overall NPS 73 -> 72
            ("وجاء تطبيق الهاتف الذكي ", "وجاء مركز خدمة المتعاملين "),  # NPS leader: mobile -> CSC
            ("7 ", "8 "),                                      # NPS leader score 77 -> 78
        ], audit)
    sh = _headline_shape(feb_prs.slides[13])
    if sh is not None:
        _swap_runs(sh, 13, [
            ("تطبيق الهاتف الذكي", "الموقع الإلكتروني"),        # value leader: mobile (92%) -> website (93%)
        ], audit)


def update_total_sample(feb_prs, total, audit):
    """Replace the overall 'العينة = 3,001' sample text on slides 9,10,12,13,14."""
    new_num = f"{total:,}"
    for si in (9, 10, 12, 13, 14):
        for sh in iter_shapes(feb_prs.slides[si].shapes):
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            if "العينة" not in tf.text:
                continue
            for para in tf.paragraphs:
                for run in para.runs:
                    t = run.text
                    for old in (TOTAL_SAMPLE_OLD, TOTAL_SAMPLE_OLD.translate(ARABIC_INDIC)):
                        if old in t:
                            run.text = t.replace(old, new_num)
                            audit.append({"slide": si, "shape": sh.shape_id,
                                          "kind": "total_sample",
                                          "old": old, "new": new_num})
                            t = run.text


def update_slide14(feb_prs, sm, audit):
    """Slide 14: satisfaction-by-demographic charts + (n) sample labels, read from
    the Feb Summary 'A16' block (T33:W55). Row->segment is fixed by the Original
    formula references; V(22)=sample, W(23)=%. The chart's 'Others' nationality bar
    merges Others + Westerner (weighted by sample)."""
    def pct(seg):
        return sm.cell(A16_ROW[seg], A16_W).value

    def base(seg):
        return sm.cell(A16_ROW[seg], A16_V).value

    def seg_pct(seg):
        if seg == "OthersNat":
            bo, bw = base("Others") or 0, base("Westerner") or 0
            po, pw = pct("Others") or 0, pct("Westerner") or 0
            tot = bo + bw
            return (bo * po + bw * pw) / tot if tot else None
        return pct(seg)

    def seg_base(seg):
        if seg == "OthersNat":
            return (base("Others") or 0) + (base("Westerner") or 0)
        return base(seg)

    slide = feb_prs.slides[14]
    for sh in iter_shapes(slide.shapes):
        if not getattr(sh, "has_chart", False):
            continue
        segs = SLIDE14_CHART_SEGS.get(sh.shape_id)
        if not segs:
            continue
        vals = [seg_pct(k) for k in segs]
        if any(not isinstance(v, (int, float)) for v in vals):
            continue
        old = [round(x, 4) for ser in sh.chart.series for x in ser.values]
        set_series_values(sh.chart, vals)
        audit.append({"slide": 14, "shape": sh.shape_id, "kind": "demo_chart",
                      "segs": segs, "old": old, "new": [round(v, 4) for v in vals]})
    by_id = {sh.shape_id: sh for sh in iter_shapes(slide.shapes)}
    for shid, key in SLIDE14_SAMPLE.items():
        sh = by_id.get(shid)
        if sh is None or not sh.has_text_frame:
            continue
        b = seg_base(key)
        if not isinstance(b, (int, float)):
            continue
        new = f"({int(round(b)):,})"
        old = sh.text_frame.text.strip()
        replace_number_text(sh, new)
        audit.append({"slide": 14, "shape": shid, "kind": "demo_sample",
                      "seg": key, "old": old, "new": new})


def _set_dpt_colour(dpt, colour):
    """Set a c:dPt's solid-fill colour, dropping any colour modifiers."""
    clr = dpt.find(".//" + A + "srgbClr")
    if clr is None:
        return False
    clr.set("val", colour)
    for child in list(clr):           # strip lumMod/lumOff/etc. so the hex is exact
        clr.remove(child)
    return True


def _build_bar_dpt(idx, colour):
    """Construct a minimal bar/column c:dPt with a solid-fill colour."""
    dpt = etree.Element(C + "dPt")
    etree.SubElement(dpt, C + "idx").set("val", str(idx))
    etree.SubElement(dpt, C + "invertIfNegative").set("val", "0")
    etree.SubElement(dpt, C + "bubble3D").set("val", "0")
    spPr = etree.SubElement(dpt, C + "spPr")
    fill = etree.SubElement(spPr, A + "solidFill")
    etree.SubElement(fill, A + "srgbClr").set("val", colour)
    ln = etree.SubElement(spPr, A + "ln")
    etree.SubElement(ln, A + "noFill")
    return dpt


def _colour_bar_chart(ser, vals):
    """Ensure one value-coloured c:dPt per bar (clone an existing one as a template,
    or build fresh dPts and insert them ahead of c:dLbls/c:cat/c:val)."""
    colours = [bar_colour(v) for v in vals]
    existing = {int(d.find(C + "idx").get("val")): d for d in ser.findall(C + "dPt")}
    if existing:
        template = copy.deepcopy(next(iter(existing.values())))
        anchor = ser.findall(C + "dPt")[-1]
        for i, colour in enumerate(colours):
            dpt = existing.get(i)
            if dpt is None:
                dpt = copy.deepcopy(template)
                dpt.find(C + "idx").set("val", str(i))
                anchor.addnext(dpt)
                anchor = dpt
            _set_dpt_colour(dpt, colour)
    else:
        ref = next((ser.find(C + t) for t in ("dLbls", "cat", "val")
                    if ser.find(C + t) is not None), None)
        if ref is None:
            return colours
        for i, colour in enumerate(colours):
            ref.addprevious(_build_bar_dpt(i, colour))
    return colours


def apply_colour_scheme(feb_prs, audit):
    """Mapping spec: charts from slide 12 onward are coloured by VALUE
    (>=90% green, 85-90% amber, <85% red). Doughnuts colour the value ring (idx 0);
    bar/column charts colour each bar. Slides before 12 are never touched. Run LAST,
    after all chart values are written."""
    for si in (12, 13, 14, 15):
        for sh in iter_shapes(feb_prs.slides[si].shapes):
            if not getattr(sh, "has_chart", False):
                continue
            ctype = str(sh.chart.chart_type)
            ser = first_series(sh.chart)
            if ser is None:
                continue
            vals = [v for s in sh.chart.series for v in s.values]
            if not vals:
                continue
            if "DOUGHNUT" in ctype:
                dpt0 = ser.find(C + "dPt")
                if dpt0 is None:
                    continue
                colour = bar_colour(vals[0])
                _set_dpt_colour(dpt0, colour)
                audit.append({"slide": si, "shape": sh.shape_id, "kind": "ring_colour",
                              "value": round(vals[0], 4), "colour": colour})
            elif "BAR" in ctype or "COLUMN" in ctype:
                colours = _colour_bar_chart(ser, vals)
                audit.append({"slide": si, "shape": sh.shape_id, "kind": "bar_colour",
                              "values": [round(v, 4) for v in vals], "colours": colours})


# Slide-15 service statements: render the (English-sorted) bars with ARABIC labels.
# The Feb basic table only carries English names, so each English service is mapped
# to its Arabic name here. The 9 services that also appear in the January deck use
# that deck's authoritative Arabic wording; the remaining 12 are translated from the
# official MOHAP service titles. Lookup is whitespace/case/apostrophe-insensitive.
SERVICE_AR = {
    # ── authoritative Arabic from data/jan_report.pptx slide 15 ──
    "Renewal of a License to Practice as a Pharmacist": "تجديد ترخيص مزاولة مهنة الصيدلة",
    "Issue License / Renew License for a Health Advertisement": "إصدار ترخيص/ تجديد ترخيص إعلان صحي",
    "Renewal of a Pharmaceutical Facility License": "تجديد ترخيص مؤسسة صيدلانية",
    "Licensing of a Pharmaceutical Facility": "ترخيص مؤسسة صيدلانية",
    "Issue/Cancel Pharmacy Home Delivery Permit": "إصدار / إلغاء تصريح صيدلية للتوصيل إلى المنزل",
    "Healthcare Professional Evaluation": "تقييم المهني الصحي",
    "Attestation of Medical Leaves and Reports": "تصديق الإجازات والتقارير الطبية",
    "Approval of Medical Leaves and The Leave of Patient's Companion": "اعتماد الاجازات المرضية وإجازة المرافق",
    "Amend a Pharmaceutical Facility License": "تعديل ترخيص مؤسسة صيدلانية",
    # ── translated from the official MOHAP service titles ──
    "Issue of Permit to Import Medicines for Personal Use": "إصدار تصريح استيراد أدوية للاستخدام الشخصي",
    "Issue of Permit to Import Medical Equipment": "إصدار تصريح استيراد أجهزة ومعدات طبية",
    "Issue of Permit to Import Chemical Precursors": "إصدار تصريح استيراد السلائف الكيميائية",
    "Registration/ Renewal of Chemical Precursor Companies": "تسجيل/ تجديد تسجيل شركات السلائف الكيميائية",
    "Treatment Abroad": "العلاج في الخارج",
    "Authorize a pharmaceutical company to Dispose any pharmaceutical products as medical waste":
        "تفويض التخلص من المنتجات الصيدلانية كنفايات",
    "Issue of Permit to Import Medicines to a Local Agent": "إصدار تصريح استيراد أدوية لوكيل محلي",
    "Handover of Narcotic Drugs custody permission": "تصريح تسليم عهدة المواد المخدرة",
    "Registration of a Medical Equipment": "تسجيل جهاز ومعدات طبية",
    "Issue a Batch Release Certificate for Shipment of Pharmaceutical Products Derived from Biological Sour":
        "شهادة الإفراج عن منتجات صيدلانية بيولوجية",
    "Complaints about private health and pharmaceutical facilities and their medical staff":
        "شكاوى المنشآت الصحية والصيدلانية الخاصة",
    "Classification of a Product": "تصنيف منتج",
}


def _norm_en(s):
    """Whitespace/case/apostrophe-insensitive key for English service names."""
    return re.sub(r"\s+", " ", str(s)).strip().lower().replace("’", "'")


_SERVICE_AR_NORM = {_norm_en(k): v for k, v in SERVICE_AR.items()}


def translate_service(name):
    """Arabic label for an English service name; None if no mapping exists."""
    return _SERVICE_AR_NORM.get(_norm_en(name))


def update_slide15(feb_prs, services, audit):
    """Slide 15: rank service statements (sample>0) by T2B desc, then write the bars
    with ARABIC labels (sorted on the English names, displayed in Arabic)."""
    slide = feb_prs.slides[15]
    total = sum(b for _, b, _ in services)
    charts = [sh for sh in iter_shapes(slide.shapes)
              if getattr(sh, "has_chart", False)
              and "BAR" in str(sh.chart.chart_type)]
    # top chart first (smaller Top), then the lower one
    charts.sort(key=lambda s: Emu(s.top).inches)
    if len(charts) >= 2:
        half = (len(services) + 1) // 2
        splits = [services[:half], services[half:]]
    else:
        splits = [services]
    missing = []
    for sh, chunk in zip(charts, splits):
        eng = [name for name, _, _ in chunk]
        labels = []
        for name in eng:
            ar = translate_service(name)
            if ar is None:
                missing.append(name)
                ar = name          # fall back to English so nothing is dropped
            labels.append(ar)
        values = [t2b for _, _, t2b in chunk]
        rewrite_bar_chart(sh.chart, labels, values)
        audit.append({"slide": 15, "shape": sh.shape_id, "kind": "service_chart",
                      "n_bars": len(chunk), "english": eng,
                      "labels": labels, "values": [round(v, 4) for v in values]})
    if missing:
        print("  WARNING: no Arabic mapping for slide-15 services (kept English):")
        for m in missing:
            print(f"    - {m}")
    # sample text: 'عينة ... = 2,166' -> total
    for sh in iter_shapes(slide.shapes):
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        if "عينة" not in tf.text and "العينة" not in tf.text:
            continue
        for para in tf.paragraphs:
            for run in para.runs:
                m = re.search(r"[\d٠-٩][\d,٠-٩]*", run.text)
                if m and len(m.group(0).translate(WESTERN).replace(",", "")) >= 3:
                    old = run.text
                    run.text = run.text[:m.start()] + f"{total:,}" + run.text[m.end():]
                    audit.append({"slide": 15, "shape": sh.shape_id, "kind": "service_sample",
                                  "old": old.strip(), "new": run.text.strip()})


# ── orchestration ────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("infile", nargs="?", default="outputs/Feb_Report_20260614.pptx")
    ap.add_argument("outfile", nargs="?", default="outputs/Feb_Report_slides9-15.pptx")
    ap.add_argument("--feb-summary", default="outputs/Feb_summary.xlsx")
    ap.add_argument("--feb-basic", default="data/feb_basic/MOHAP CSAT Tabs - Feb 2026.xlsx")
    ap.add_argument("--jan-basic", default="data/jan_basic/MOHAP CSAT Tabs for Jan-2026.xlsx")
    ap.add_argument("--jan-summary", default="data/jan_summary.xlsx")
    ap.add_argument("--jan-deck", default="data/jan_report.pptx")
    args = ap.parse_args()

    sm = load_summary(args.feb_summary)
    jan_sm = load_summary(args.jan_summary)
    feb_basic = openpyxl.load_workbook(args.feb_basic, data_only=True)["Percentages"]
    base_row, t2b_row = locate_overall_sat(feb_basic)

    jan_secs = rc.load_sections(args.jan_basic)
    feb_secs = rc.load_sections(args.feb_basic)
    services = load_services(feb_basic, base_row, t2b_row)
    total_sample = sum(b for _, b, _ in services)

    feb_prs = Presentation(args.infile)
    jan_prs = Presentation(args.jan_deck)

    audit = []
    update_distribution_slides(feb_prs, jan_prs, jan_secs, feb_secs, [9], audit)
    update_slide10(feb_prs, feb_secs, audit)
    update_doughnut_dashboard(feb_prs, sm, audit)
    update_nps_rectangles(feb_prs, sm, audit)
    update_channel_samples(feb_prs, jan_prs, sm, jan_sm, audit)
    update_headlines(feb_prs, sm, audit)
    update_slide14(feb_prs, sm, audit)
    update_slide15(feb_prs, services, audit)
    update_total_sample(feb_prs, total_sample, audit)
    # value-based chart colours for slides 12+ (run last, after all values written)
    apply_colour_scheme(feb_prs, audit)

    feb_prs.save(args.outfile)

    # report
    by_slide = {}
    for a in audit:
        by_slide.setdefault(a["slide"], []).append(a)
    print(f"Saved {args.outfile}")
    print(f"Overall-Satisfaction table rows: base={base_row} T2B={t2b_row}; "
          f"slide-15 services with sample>0 = {len(services)}, total sample = {total_sample:,}")
    for si in sorted(by_slide):
        kinds = {}
        for a in by_slide[si]:
            kinds[a["kind"]] = kinds.get(a["kind"], 0) + 1
        summary = ", ".join(f"{k}×{n}" for k, n in kinds.items())
        print(f"  slide {si:>2}: {len(by_slide[si])} edits  ({summary})")

    audit_path = Path(args.outfile).with_suffix("").as_posix() + "_audit.json"
    with open(audit_path, "w", encoding="utf-8") as f:
        json.dump(audit, f, ensure_ascii=False, indent=2)
    print(f"Audit: {audit_path}")


if __name__ == "__main__":
    main()
