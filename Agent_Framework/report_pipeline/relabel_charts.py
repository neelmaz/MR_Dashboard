"""
Label-anchored chart updater for the distribution / demographic-profile charts.

The value-matching chart resolver (update_charts.py) picks the wrong source cell
when a percentage appears in many places (e.g. the Gender split resolved to
0.448/0.552 instead of the real 0.457/0.543). This writer instead anchors on the
human-readable LABELS, which are stable across months:

  1. Enumerate the workbook's "distribution sections" — a header row (e.g.
     "Gender", "ID3. Which age group...") followed by "Base: All respondents",
     then labelled option rows with a Total-column percentage, terminated by a
     "Total = 1" row. The 12 sections align 1:1 between Jan and Feb by header.
  2. For each chart series in the JANUARY deck, match its value vector to exactly
     one section, recording per-bar which option label(s) feed it (a bar may be
     the SUM of several option rows, e.g. nationality "Others" = Westerner+Others).
  3. Read the FEBRUARY value for each bar by the same section header + label(s),
     and write it into the corresponding chart of the February deck (charts are
     matched 1:1 by (slide index, shape id) — the Feb deck is a copy of Jan).

Charts that do not map cleanly to a single section are left untouched and listed
in the audit, so nothing outside the distribution charts is disturbed.

Usage:
  python relabel_charts.py <feb_in.pptx> <feb_out.pptx>
    [--jan-deck data/jan_report.pptx]
    [--jan-xlsx ...] [--feb-xlsx ...]
"""

from __future__ import annotations

import argparse
import json
import logging
from collections import defaultdict
from itertools import combinations
from pathlib import Path

import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger("relabel_charts")

C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
NS = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}

TOL = 0.0015           # percentage points stored to 3 dp
TOTAL_COL = 2          # column B holds the "Total" banner
JAN_XLSX = "data/jan_basic/MOHAP CSAT Tabs for Jan-2026.xlsx"
FEB_XLSX = "data/feb_basic/MOHAP CSAT Tabs - Feb 2026.xlsx"
JAN_DECK = "data/jan_report.pptx"


def _num(v):
    return v if isinstance(v, (int, float)) and not isinstance(v, bool) else None


def load_sections(path):
    """{header: {label: total_value}} for every distribution section."""
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb["Percentages"]
    nrow = ws.max_row
    secs = {}
    r = 1
    while r <= nrow:
        a = ws.cell(r, 1).value
        nxt = ws.cell(r + 1, 1).value if r + 1 <= nrow else None
        if (isinstance(a, str) and a.strip()
                and isinstance(nxt, str) and "base: all resp" in nxt.lower()):
            header = a.strip()
            labels = {}
            rr = r + 2
            while rr <= nrow:
                lab = ws.cell(rr, 1).value
                if not isinstance(lab, str):
                    break
                ls = lab.strip()
                if ls.lower() == "total" or ls == "_____" or ls == "":
                    break
                labels[ls] = _num(ws.cell(rr, TOTAL_COL).value)
                rr += 1
            secs[header] = labels
        r += 1
    return secs


# ── chart traversal ─────────────────────────────────────────────────────────────

def iter_charts(shapes):
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_charts(s.shapes)
        elif getattr(s, "has_chart", False):
            yield s


def series_pts(ser_el):
    """Ordered [(idx, c:v element, float)] of a series' value cache."""
    out = []
    val = ser_el.find(f"{C}val")
    if val is None:
        return out
    for pt in val.iter(f"{C}pt"):
        v = pt.find(f"{C}v")
        if v is None or v.text is None:
            continue
        try:
            f = float(v.text)
        except ValueError:
            continue
        out.append((int(pt.get("idx")), v, f))
    out.sort(key=lambda x: x[0])
    return out


# ── matching a series to a section ──────────────────────────────────────────────

def match_series_to_section(values, sec_labels):
    """Try to explain `values` (in order) using the section's {label: value}.

    Returns a list (one per position) of label-lists whose Jan values sum to the
    bar, or None if the section can't account for every bar. Each label is used
    at most once across the whole series.

    A match is only accepted when it is a genuine *full distribution*: at least
    two bars, and every section option carrying material mass (> TOL) is consumed
    by some bar. This rejects the coincidental single-value / doughnut matches
    (e.g. a lone 0.013 or a [v, 1-v] KPI) that merely happen to equal one option.
    """
    if len(values) < 2:
        return None
    avail = {lab: v for lab, v in sec_labels.items() if v is not None}
    recipe = [None] * len(values)

    # pass 1: unique single-label exact matches
    for i, v in enumerate(values):
        hits = [lab for lab, lv in avail.items() if abs(lv - v) <= TOL]
        if len(hits) == 1:
            recipe[i] = [hits[0]]
            del avail[hits[0]]

    # pass 2: remaining single-label matches (pick any exact)
    for i, v in enumerate(values):
        if recipe[i] is not None:
            continue
        hit = next((lab for lab, lv in avail.items() if abs(lv - v) <= TOL), None)
        if hit is not None:
            recipe[i] = [hit]
            del avail[hit]

    # pass 3: merged buckets (sum of 2..3 leftover labels)
    for i, v in enumerate(values):
        if recipe[i] is not None:
            continue
        found = None
        items = list(avail.items())
        for k in (2, 3):
            for combo in combinations(items, k):
                if abs(sum(lv for _, lv in combo) - v) <= TOL:
                    found = [lab for lab, _ in combo]
                    break
            if found:
                break
        if found:
            recipe[i] = found
            for lab in found:
                del avail[lab]

    if any(r is None for r in recipe):
        return None
    # full-distribution gate: no material option may be left unconsumed
    if any(v is not None and v > TOL for v in avail.values()):
        return None
    # one-bar-per-option gate: a true distribution chart has roughly as many bars
    # as the section has material options (at most one tail bucket merges two
    # small options, e.g. nationality "Others" = Westerner + Others). This rejects
    # [v, 1-v] doughnut KPIs that "explain" a 6-option section by merging it.
    material = sum(1 for v in sec_labels.values() if v is not None and v > TOL)
    if len(values) < material - 1:
        return None
    return recipe


def best_section(values, sections):
    """Return (header, recipe) for the section that explains the series, or None.
    Prefers the section using the fewest leftover labels (tightest fit)."""
    best = None
    for header, labels in sections.items():
        recipe = match_series_to_section(values, labels)
        if recipe is None:
            continue
        used = sum(len(r) for r in recipe)
        leftover = len([v for v in labels.values() if v is not None]) - used
        score = (leftover, len(header))
        if best is None or score < best[0]:
            best = (score, header, recipe)
    return (best[1], best[2]) if best else None


def feb_value(feb_sections, header, labels):
    vals = feb_sections.get(header, {})
    tot = 0.0
    for lab in labels:
        v = vals.get(lab)
        if v is None:
            return None
        tot += v
    return tot


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("feb_in")
    ap.add_argument("feb_out")
    ap.add_argument("--jan-deck", default=JAN_DECK)
    ap.add_argument("--jan-xlsx", default=JAN_XLSX)
    ap.add_argument("--feb-xlsx", default=FEB_XLSX)
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    jan_sections = load_sections(args.jan_xlsx)
    feb_sections = load_sections(args.feb_xlsx)
    logger.info(f"sections: jan={len(jan_sections)} feb={len(feb_sections)}")

    jan_prs = Presentation(args.jan_deck)
    feb_prs = Presentation(args.feb_in)

    # index Jan charts by (slide, shape_id) -> chart element
    jan_charts = {}
    for si, slide in enumerate(jan_prs.slides):
        for sh in iter_charts(slide.shapes):
            jan_charts[(si, sh.shape_id)] = sh.chart

    audit = []
    matched = changed = 0
    for si, slide in enumerate(feb_prs.slides):
        for sh in iter_charts(slide.shapes):
            jchart = jan_charts.get((si, sh.shape_id))
            if jchart is None:
                continue
            jcs = jchart._chartSpace
            fcs = sh.chart._chartSpace
            jsers = jcs.findall(".//c:ser", NS)
            fsers = fcs.findall(".//c:ser", NS)
            for jser, fser in zip(jsers, fsers):
                jpts = series_pts(jser)
                fpts = series_pts(fser)
                if not jpts or len(jpts) != len(fpts):
                    continue
                jvals = [f for _, _, f in jpts]
                m = best_section(jvals, jan_sections)
                if m is None:
                    continue
                header, recipe = m
                feb = [feb_value(feb_sections, header, labs) for labs in recipe]
                if any(v is None for v in feb):
                    audit.append({"slide": si, "shape_id": sh.shape_id,
                                  "section": header, "status": "feb_missing",
                                  "recipe": recipe})
                    continue
                matched += 1
                old = [f for _, _, f in fpts]
                diff = any(abs(o - n) > TOL for o, n in zip(old, feb))
                if diff and not args.dry_run:
                    for (_, vel, _), nv in zip(fpts, feb):
                        vel.text = f"{nv:.10g}"
                if diff:
                    changed += 1
                audit.append({
                    "slide": si, "shape_id": sh.shape_id, "section": header,
                    "status": "changed" if diff else "already_correct",
                    "recipe": recipe,
                    "jan": [round(v, 4) for v in jvals],
                    "old_feb": [round(v, 4) for v in old],
                    "new_feb": [round(v, 4) for v in feb],
                })

    if not args.dry_run:
        feb_prs.save(args.feb_out)
    audit_path = Path(args.feb_out).with_name(Path(args.feb_out).stem + "_relabel_audit.json")
    with open(audit_path, "w", encoding="utf-8") as f:
        json.dump({"matched_series": matched, "changed_series": changed,
                   "entries": audit}, f, ensure_ascii=False, indent=2)
    logger.info(f"matched series: {matched}  changed: {changed}")
    logger.info(f"audit: {audit_path}")
    if args.dry_run:
        logger.info("DRY RUN — no deck written")


if __name__ == "__main__":
    main()
